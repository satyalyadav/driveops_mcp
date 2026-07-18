"""Google Drive API adapter used by DriveOps tools."""

from __future__ import annotations

import base64
import io
import mimetypes
import re
import shutil
import stat
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload, MediaIoBaseUpload

from .auth import get_credentials
from .schemas import (
    GOOGLE_DOC_MIME,
    GOOGLE_FOLDER_MIME,
    GOOGLE_SHEET_MIME,
    GOOGLE_SLIDE_MIME,
    normalize_file,
)


class DriveOpsError(RuntimeError):
    pass


class AmbiguousFileError(DriveOpsError):
    def __init__(
        self, file_name: str, matches: list[dict[str, Any]], has_more: bool = False
    ) -> None:
        super().__init__(f"Multiple files named '{file_name}' found.")
        self.file_name = file_name
        self.matches = matches
        self.has_more = has_more


class GoogleDriveClient:
    def __init__(
        self, drive_service: Any | None = None, sheets_service: Any | None = None
    ) -> None:
        if drive_service is None:
            creds = get_credentials()
            drive_service = build("drive", "v3", credentials=creds)
            sheets_service = build("sheets", "v4", credentials=creds)
        self.drive = drive_service
        self.sheets = sheets_service

    @staticmethod
    def _escape(value: str) -> str:
        return value.replace("'", "\\'")

    @staticmethod
    def _looks_like_id(value: str) -> bool:
        return bool(re.fullmatch(r"[A-Za-z0-9_-]{10,}", value or ""))

    def _execute_files_page(
        self, *, max_items: int | None = None, **params: Any
    ) -> tuple[list[dict[str, Any]], str | None]:
        items: list[dict[str, Any]] = []
        page_token = params.pop("pageToken", None)
        while True:
            call_params = dict(params)
            if max_items is not None:
                remaining = max_items - len(items)
                if remaining <= 0:
                    break
                requested = int(call_params.get("pageSize", remaining))
                call_params["pageSize"] = max(1, min(requested, remaining))
            if page_token:
                call_params["pageToken"] = page_token
            resp = self.drive.files().list(**call_params).execute()
            items.extend(resp.get("files", []))
            page_token = resp.get("nextPageToken")
            if max_items is not None and len(items) >= max_items:
                break
            if not page_token:
                break
        return items[:max_items] if max_items is not None else items, page_token

    def _execute_files_list(self, **params: Any) -> list[dict[str, Any]]:
        files, _ = self._execute_files_page(**params)
        return files

    def get_file(self, file_id: str, fields: str | None = None) -> dict[str, Any]:
        fields = (
            fields
            or "id,name,mimeType,createdTime,modifiedTime,size,parents,webViewLink,webContentLink"
        )
        return (
            self.drive.files()
            .get(
                fileId=file_id,
                fields=fields,
                supportsAllDrives=True,
            )
            .execute()
        )

    def resolve_file(self, file_id_or_name: str) -> dict[str, Any]:
        value = file_id_or_name.strip()
        if not value:
            raise DriveOpsError("file_id_or_name is required.")
        if self._looks_like_id(value):
            try:
                return self.get_file(value)
            except Exception:
                pass

        exact, has_more = self._exact_file_matches(value)
        if len(exact) > 1:
            raise AmbiguousFileError(value, self.enrich_files(exact), has_more)
        if exact:
            return exact[0]
        search = self.search_files(query=value, page_size=1)
        files = search.get("files", [])
        if files:
            return files[0]
        raise DriveOpsError(f"File '{file_id_or_name}' not found.")

    def _exact_file_matches(
        self, value: str, max_matches: int = 6
    ) -> tuple[list[dict[str, Any]], bool]:
        safe = self._escape(value)
        return self._execute_files_page(
            max_items=max_matches,
            q=f"name = '{safe}' and trashed = false",
            pageSize=max_matches,
            fields="files(id,name,mimeType,createdTime,modifiedTime,size,parents,webViewLink,webContentLink)",
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            corpora="allDrives",
        )

    def resolve_folder(self, folder_id_or_name: str) -> dict[str, Any]:
        value = folder_id_or_name.strip()
        if not value:
            raise DriveOpsError("folder_id_or_name is required.")
        if value.lower() in {"root", "my drive", "drive", "/"}:
            return {
                "id": "root",
                "name": "My Drive",
                "mimeType": GOOGLE_FOLDER_MIME,
                "parents": [],
            }
        if self._looks_like_id(value):
            folder = self.get_file(value, fields="id,name,mimeType,parents")
            if folder.get("mimeType") != GOOGLE_FOLDER_MIME:
                raise DriveOpsError(f"{value} is not a Google Drive folder.")
            return folder

        cleaned = value
        if cleaned.lower().startswith("my ") and cleaned.lower() != "my drive":
            cleaned = cleaned[3:].strip()
        if cleaned.lower().endswith(" folder"):
            cleaned = cleaned[:-7].strip()
        safe = self._escape(cleaned)
        candidates = self._execute_files_list(
            q=f"name contains '{safe}' and mimeType = '{GOOGLE_FOLDER_MIME}' and trashed = false",
            pageSize=10,
            fields="files(id,name,mimeType,parents)",
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            corpora="allDrives",
        )
        for item in candidates:
            if item.get("name", "").lower() == cleaned.lower():
                return item
        if candidates:
            return candidates[0]
        raise DriveOpsError(f"Folder '{folder_id_or_name}' not found.")

    def folder_path(self, parent_id: str, cache: dict[str, str] | None = None) -> str:
        cache = cache if cache is not None else {}
        if not parent_id or parent_id == "root":
            return "My Drive"
        if parent_id in cache:
            return cache[parent_id]
        try:
            folder = self.get_file(parent_id, fields="id,name,parents")
            parents = folder.get("parents", [])
            name = folder.get("name", "Unknown")
            if not parents or parents[0] == "root":
                path = f"My Drive > {name}" if name != "My Drive" else "My Drive"
            else:
                path = f"{self.folder_path(parents[0], cache)} > {name}"
            cache[parent_id] = path
            return path
        except Exception:
            return "Unknown"

    def enrich_files(self, files: list[dict[str, Any]]) -> list[dict[str, Any]]:
        cache: dict[str, str] = {}
        enriched = []
        for item in files:
            parents = item.get("parents", [])
            paths = [self.folder_path(parent, cache) for parent in parents] or [
                "My Drive"
            ]
            item = dict(item)
            item["folderPaths"] = paths
            item["folderPath"] = paths[0]
            if not item.get("webViewLink") and item.get("id"):
                item["webViewLink"] = (
                    f"https://drive.google.com/file/d/{item['id']}/view"
                )
            enriched.append(normalize_file(item))
        return enriched

    def search_files(
        self,
        *,
        query: str,
        folder_id: str | None = None,
        mime_types: list[str] | None = None,
        page_size: int = 10,
        page_token: str | None = None,
    ) -> dict[str, Any]:
        query = (query or "*").strip()
        page_size = max(1, min(int(page_size), 100))
        clauses = ["trashed = false"]
        if folder_id:
            folder = self.resolve_folder(folder_id)
            clauses.append(f"'{folder['id']}' in parents")
        if mime_types:
            mime_clause = " or ".join(
                f"mimeType = '{self._escape(m)}'" for m in mime_types
            )
            clauses.append(f"({mime_clause})")

        raw_query = any(
            op in query for op in [" contains ", "=", " in parents", "fullText"]
        )
        if query != "*":
            if raw_query:
                clauses.insert(0, f"({query})")
            else:
                tokens = [t for t in re.split(r"\s+", query) if t]
                token_clauses = []
                for token in tokens:
                    safe = self._escape(token)
                    token_clauses.append(
                        f"(name contains '{safe}' or fullText contains '{safe}')"
                    )
                if token_clauses:
                    clauses.insert(0, "(" + " and ".join(token_clauses) + ")")

        files, next_page_token = self._execute_files_page(
            max_items=page_size,
            q=" and ".join(clauses),
            pageSize=page_size,
            fields="nextPageToken, files(id,name,mimeType,createdTime,modifiedTime,size,parents,webViewLink,webContentLink)",
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            corpora="allDrives",
            pageToken=page_token,
        )
        return {
            "query": query,
            "count": len(files),
            "has_more": bool(next_page_token),
            "next_page_token": next_page_token,
            "files": self.enrich_files(files),
        }

    def list_folder(
        self,
        folder_id_or_name: str,
        page_size: int = 100,
        page_token: str | None = None,
    ) -> dict[str, Any]:
        folder = self.resolve_folder(folder_id_or_name)
        page_size = max(1, min(int(page_size), 1000))
        files, next_page_token = self._execute_files_page(
            max_items=page_size,
            q=f"'{folder['id']}' in parents and trashed = false",
            pageSize=page_size,
            fields="nextPageToken, files(id,name,mimeType,createdTime,modifiedTime,size,parents,webViewLink,webContentLink)",
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            corpora="allDrives",
            pageToken=page_token,
        )
        return {
            "folder": folder,
            "count": len(files),
            "has_more": bool(next_page_token),
            "next_page_token": next_page_token,
            "files": self.enrich_files(files),
        }

    def get_changes(
        self,
        folder_id_or_name: str,
        since: str,
        page_size: int = 100,
        page_token: str | None = None,
    ) -> dict[str, Any]:
        folder = self.resolve_folder(folder_id_or_name)
        try:
            date_part = since.split("T", 1)[0]
            since_dt = datetime.fromisoformat(date_part)
        except ValueError as exc:
            raise DriveOpsError("since must be an ISO date like YYYY-MM-DD.") from exc
        since_str = since_dt.strftime("%Y-%m-%dT00:00:00")
        page_size = max(1, min(int(page_size), 1000))
        files, next_page_token = self._execute_files_page(
            max_items=page_size,
            q=f"'{folder['id']}' in parents and modifiedTime >= '{since_str}' and trashed = false",
            pageSize=page_size,
            fields="nextPageToken, files(id,name,mimeType,createdTime,modifiedTime,size,parents,webViewLink,webContentLink)",
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            corpora="allDrives",
            orderBy="modifiedTime desc",
            pageToken=page_token,
        )
        enriched = self.enrich_files(files)
        for item in enriched:
            item["changeType"] = (
                "new" if (item.get("createdTime") or "") >= since_str else "edited"
            )
        return {
            "folder": folder,
            "since": since,
            "count": len(enriched),
            "has_more": bool(next_page_token),
            "next_page_token": next_page_token,
            "changedFiles": enriched,
        }

    def read_file(
        self,
        file_id_or_name: str,
        export_format: str | None = None,
        max_text_bytes: int = 200_000,
    ) -> dict[str, Any]:
        try:
            resolved = self.resolve_file(file_id_or_name)
        except AmbiguousFileError as exc:
            return {
                "status": "ambiguous",
                "contentType": "ambiguous",
                "text": None,
                "name": exc.file_name,
                "matches": exc.matches,
                "has_more": exc.has_more,
                "message": "Multiple files matched this name. Ask the user to choose one by path, modified time, or file ID.",
            }
        file_id = resolved["id"]
        meta = self.get_file(
            file_id, fields="id,name,mimeType,size,webViewLink,webContentLink"
        )
        mime = meta.get("mimeType", "")
        export_mime = export_format
        if mime == GOOGLE_DOC_MIME:
            export_mime = export_mime or "text/plain"
        elif mime == GOOGLE_SHEET_MIME:
            export_mime = export_mime or "text/csv"
        elif mime == GOOGLE_SLIDE_MIME:
            export_mime = export_mime or "text/plain"

        if mime.startswith("application/vnd.google-apps"):
            request = self.drive.files().export(
                fileId=file_id, mimeType=export_mime or "application/pdf"
            )
            data = self._download_request(request, max_text_bytes + 1)
            return self._content_response(
                meta, data, export_mime or "application/octet-stream", max_text_bytes
            )

        if mime.startswith("text/") or mime in {"application/json", "application/xml"}:
            request = self.drive.files().get_media(fileId=file_id)
            data = self._download_request(request, max_text_bytes + 1)
            return self._content_response(meta, data, mime, max_text_bytes)

        return {
            "id": file_id,
            "name": meta.get("name"),
            "mimeType": mime,
            "contentType": "download_hint",
            "text": None,
            "truncated": False,
            "webViewLink": meta.get("webViewLink"),
            "webContentLink": meta.get("webContentLink"),
            "message": "Binary or large non-text content was not loaded into model context.",
        }

    def download_file(
        self,
        file_id_or_name: str,
        *,
        export_format: str | None = None,
        output_path: str | None = None,
        max_bytes: int = 25_000_000,
        overwrite: bool = False,
    ) -> dict[str, Any]:
        """Download a blob or export a Workspace file as actual bytes.

        With ``output_path`` the bytes are written on the MCP server host. Without
        it, reasonably sized content is returned as base64 for remote MCP clients.
        """

        meta = self.resolve_file(file_id_or_name)
        meta = self.get_file(
            meta["id"], fields="id,name,mimeType,size,webViewLink,webContentLink"
        )
        mime = meta.get("mimeType", "")
        content_type = mime
        if mime.startswith("application/vnd.google-apps"):
            export_format = export_format or self._default_export_format(mime)
            request = self.drive.files().export(
                fileId=meta["id"], mimeType=export_format
            )
            content_type = export_format
        else:
            request = self.drive.files().get_media(
                fileId=meta["id"], supportsAllDrives=True
            )

        max_bytes = max(1, min(int(max_bytes), 100_000_000))
        data = self._download_request(request, max_bytes + 1)
        if len(data) > max_bytes:
            raise DriveOpsError(
                f"Download exceeds the {max_bytes}-byte safety limit. Increase max_bytes up to 100 MB."
            )

        filename = self._download_name(meta.get("name") or meta["id"], content_type)
        result = {
            "id": meta["id"],
            "name": meta.get("name"),
            "filename": filename,
            "mimeType": mime,
            "contentType": content_type,
            "size": len(data),
            "webViewLink": meta.get("webViewLink"),
        }
        if output_path:
            destination = Path(output_path).expanduser()
            if destination.exists() and destination.is_dir():
                destination = destination / filename
            if destination.exists() and not overwrite:
                raise DriveOpsError(
                    f"Download destination already exists: {destination}. Set overwrite=true to replace it."
                )
            destination.parent.mkdir(parents=True, exist_ok=True)
            destination.write_bytes(data)
            result["saved_path"] = str(destination.resolve())
        else:
            result["content_base64"] = base64.b64encode(data).decode("ascii")
        return result

    def extract_file(
        self,
        file_id_or_name: str,
        *,
        output_dir: str | None = None,
        max_bytes: int = 50_000_000,
        max_text_chars: int = 500_000,
        overwrite: bool = False,
    ) -> dict[str, Any]:
        """Extract readable text or safely unpack a ZIP archive."""

        resolved = self.resolve_file(file_id_or_name)
        meta = self.get_file(resolved["id"], fields="id,name,mimeType,size,webViewLink")
        mime = meta.get("mimeType", "")
        export_format = None
        if mime == GOOGLE_DOC_MIME:
            export_format = "text/plain"
        elif mime == GOOGLE_SHEET_MIME:
            export_format = "text/csv"
        elif mime == GOOGLE_SLIDE_MIME:
            export_format = "text/plain"

        if mime.startswith("application/vnd.google-apps"):
            if not export_format:
                raise DriveOpsError(f"Text extraction is not supported for {mime}.")
            request = self.drive.files().export(
                fileId=meta["id"], mimeType=export_format
            )
        else:
            request = self.drive.files().get_media(
                fileId=meta["id"], supportsAllDrives=True
            )
        max_bytes = max(1, min(int(max_bytes), 100_000_000))
        data = self._download_request(request, max_bytes + 1)
        if len(data) > max_bytes:
            raise DriveOpsError(
                f"File exceeds the {max_bytes}-byte extraction safety limit."
            )

        name = (meta.get("name") or "").lower()
        result: dict[str, Any] = {
            "id": meta["id"],
            "name": meta.get("name"),
            "mimeType": mime,
            "source_size": len(data),
            "webViewLink": meta.get("webViewLink"),
        }
        if (
            mime.startswith("text/")
            or mime in {"application/json", "application/xml"}
            or export_format
        ):
            text = data.decode("utf-8", errors="replace")
        elif mime == "application/pdf" or name.endswith(".pdf"):
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
            result["pages"] = len(reader.pages)
        elif (
            mime
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or name.endswith(".docx")
        ):
            from docx import Document

            document = Document(io.BytesIO(data))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        elif (
            mime
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or name.endswith(".pptx")
        ):
            from pptx import Presentation

            presentation = Presentation(io.BytesIO(data))
            slide_text = []
            for number, slide in enumerate(presentation.slides, start=1):
                parts = [
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                ]
                slide_text.append(f"Slide {number}\n" + "\n".join(parts))
            text = "\n\n".join(slide_text)
            result["slides"] = len(presentation.slides)
        elif (
            mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            or name.endswith(".xlsx")
        ):
            from openpyxl import load_workbook

            workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            sheets = []
            for sheet in workbook.worksheets:
                rows = [
                    "\t".join("" if value is None else str(value) for value in row)
                    for row in sheet.iter_rows(values_only=True)
                ]
                sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
            text = "\n\n".join(sheets)
            result["sheets"] = workbook.sheetnames
        elif mime in {
            "application/zip",
            "application/x-zip-compressed",
        } or name.endswith(".zip"):
            return self._extract_zip(
                meta, data, output_dir=output_dir, overwrite=overwrite
            )
        else:
            raise DriveOpsError(
                "Unsupported extraction format. Supported: Google Docs/Sheets/Slides, text, PDF, DOCX, PPTX, XLSX, and ZIP."
            )

        max_text_chars = max(1, min(int(max_text_chars), 2_000_000))
        result.update(
            {
                "contentType": "text",
                "text": text[:max_text_chars],
                "truncated": len(text) > max_text_chars,
                "characters": len(text),
            }
        )
        return result

    @staticmethod
    def _extract_zip(
        meta: dict[str, Any], data: bytes, *, output_dir: str | None, overwrite: bool
    ) -> dict[str, Any]:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = archive.infolist()
            if len(entries) > 1000:
                raise DriveOpsError("ZIP archive contains more than 1000 entries.")
            expanded_size = sum(entry.file_size for entry in entries)
            if expanded_size > 100_000_000:
                raise DriveOpsError(
                    "ZIP archive expands beyond the 100 MB safety limit."
                )
            manifest = [
                {
                    "name": entry.filename,
                    "size": entry.file_size,
                    "is_dir": entry.is_dir(),
                }
                for entry in entries
            ]
            result: dict[str, Any] = {
                "id": meta["id"],
                "name": meta.get("name"),
                "mimeType": meta.get("mimeType"),
                "contentType": "archive",
                "entry_count": len(entries),
                "expanded_size": expanded_size,
                "entries": manifest,
            }
            if not output_dir:
                result["message"] = (
                    "Archive inspected but not unpacked; provide output_dir to extract it."
                )
                return result

            root = Path(output_dir).expanduser().resolve()
            root.mkdir(parents=True, exist_ok=True)
            destinations: list[tuple[zipfile.ZipInfo, Path]] = []
            for entry in entries:
                mode = entry.external_attr >> 16
                if stat.S_ISLNK(mode):
                    raise DriveOpsError(
                        f"Refusing to extract ZIP symlink: {entry.filename}"
                    )
                destination = (root / entry.filename).resolve()
                if not destination.is_relative_to(root):
                    raise DriveOpsError(f"Refusing unsafe ZIP path: {entry.filename}")
                if not entry.is_dir() and destination.exists() and not overwrite:
                    raise DriveOpsError(
                        f"Extraction destination already exists: {destination}"
                    )
                destinations.append((entry, destination))
            extracted = []
            for entry, destination in destinations:
                if entry.is_dir():
                    destination.mkdir(parents=True, exist_ok=True)
                    continue
                destination.parent.mkdir(parents=True, exist_ok=True)
                with archive.open(entry) as source, destination.open("wb") as target:
                    shutil.copyfileobj(source, target)
                extracted.append(str(destination))
            result.update({"output_dir": str(root), "extracted_paths": extracted})
            return result

    @staticmethod
    def _default_export_format(mime_type: str) -> str:
        formats = {
            GOOGLE_DOC_MIME: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            GOOGLE_SHEET_MIME: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            GOOGLE_SLIDE_MIME: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.google-apps.drawing": "application/pdf",
        }
        if mime_type not in formats:
            raise DriveOpsError(
                f"No default export format is known for {mime_type}; provide export_format."
            )
        return formats[mime_type]

    @staticmethod
    def _download_name(name: str, content_type: str) -> str:
        suffix = mimetypes.guess_extension(content_type) or ""
        if suffix == ".jpe":
            suffix = ".jpg"
        if suffix and not name.lower().endswith(suffix.lower()):
            return f"{name}{suffix}"
        return name

    @staticmethod
    def _download_request(request: Any, max_bytes: int) -> bytes:
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done and fh.tell() <= max_bytes:
            _, done = downloader.next_chunk()
        return fh.getvalue()

    @staticmethod
    def _content_response(
        meta: dict[str, Any], data: bytes, content_type: str, max_text_bytes: int
    ) -> dict[str, Any]:
        truncated = len(data) > max_text_bytes
        raw = data[:max_text_bytes]
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("utf-8", errors="replace")
        return {
            "id": meta.get("id"),
            "name": meta.get("name"),
            "mimeType": meta.get("mimeType"),
            "contentType": content_type,
            "text": text,
            "truncated": truncated,
            "webViewLink": meta.get("webViewLink"),
            "webContentLink": meta.get("webContentLink"),
        }

    def create_folder(self, name: str, parent_id: str) -> dict[str, Any]:
        body = {"name": name, "mimeType": GOOGLE_FOLDER_MIME, "parents": [parent_id]}
        return (
            self.drive.files()
            .create(
                body=body,
                fields="id,name,mimeType,parents",
                supportsAllDrives=True,
            )
            .execute()
        )

    def create_file(
        self,
        *,
        name: str,
        parent_id: str,
        mime_type: str | None = None,
        text: str | None = None,
        content_base64: str | None = None,
        local_path: str | None = None,
    ) -> dict[str, Any]:
        sources = sum(value is not None for value in (text, content_base64, local_path))
        if sources > 1:
            raise DriveOpsError(
                "Provide only one of text, content_base64, or local_path."
            )
        body = {"name": name, "parents": [parent_id]}
        media = None
        if local_path is not None:
            path = Path(local_path).expanduser()
            if not path.is_file():
                raise DriveOpsError(f"Upload source is not a file: {path}")
            mime_type = (
                mime_type
                or mimetypes.guess_type(path.name)[0]
                or "application/octet-stream"
            )
            media = MediaFileUpload(str(path), mimetype=mime_type, resumable=True)
        elif content_base64 is not None:
            try:
                raw = base64.b64decode(content_base64, validate=True)
            except ValueError as exc:
                raise DriveOpsError("content_base64 is not valid base64.") from exc
            mime_type = mime_type or "application/octet-stream"
            media = MediaIoBaseUpload(
                io.BytesIO(raw), mimetype=mime_type, resumable=True
            )
        elif text is not None:
            mime_type = mime_type or "text/plain"
            media = MediaIoBaseUpload(
                io.BytesIO(text.encode("utf-8")), mimetype=mime_type, resumable=True
            )
        elif mime_type:
            body["mimeType"] = mime_type

        return (
            self.drive.files()
            .create(
                body=body,
                media_body=media,
                fields="id,name,mimeType,parents,createdTime,modifiedTime,size,webViewLink,webContentLink",
                supportsAllDrives=True,
            )
            .execute()
        )

    def rename_file(self, file_id: str, new_name: str) -> dict[str, Any]:
        return (
            self.drive.files()
            .update(
                fileId=file_id,
                body={"name": new_name},
                fields="id,name,mimeType,parents,modifiedTime,webViewLink",
                supportsAllDrives=True,
            )
            .execute()
        )

    def copy_file(
        self, file_id: str, *, name: str | None = None, parent_id: str | None = None
    ) -> dict[str, Any]:
        body: dict[str, Any] = {}
        if name:
            body["name"] = name
        if parent_id:
            body["parents"] = [parent_id]
        return (
            self.drive.files()
            .copy(
                fileId=file_id,
                body=body,
                fields="id,name,mimeType,parents,createdTime,modifiedTime,size,webViewLink,webContentLink",
                supportsAllDrives=True,
            )
            .execute()
        )

    def move_file(
        self, file_id: str, add_parent: str, remove_parent: str
    ) -> dict[str, Any]:
        return (
            self.drive.files()
            .update(
                fileId=file_id,
                addParents=add_parent,
                removeParents=remove_parent,
                fields="id,name,mimeType,parents,modifiedTime",
                supportsAllDrives=True,
            )
            .execute()
        )

    def set_trashed(self, file_id: str, trashed: bool) -> dict[str, Any]:
        return (
            self.drive.files()
            .update(
                fileId=file_id,
                body={"trashed": bool(trashed)},
                fields="id,name,mimeType,parents,trashed,modifiedTime",
                supportsAllDrives=True,
            )
            .execute()
        )

    def delete_file(self, file_id: str) -> dict[str, Any]:
        self.drive.files().delete(fileId=file_id, supportsAllDrives=True).execute()
        return {"id": file_id, "deleted": True}

    def list_permissions(
        self,
        file_id_or_name: str,
        *,
        page_size: int = 100,
        page_token: str | None = None,
    ) -> dict[str, Any]:
        file = self.resolve_file(file_id_or_name)
        params: dict[str, Any] = {
            "fileId": file["id"],
            "pageSize": max(1, min(int(page_size), 100)),
            "fields": "nextPageToken,permissions(id,type,role,emailAddress,domain,displayName,expirationTime,deleted,allowFileDiscovery,permissionDetails)",
            "supportsAllDrives": True,
        }
        if page_token:
            params["pageToken"] = page_token
        response = self.drive.permissions().list(**params).execute()
        return {
            "file": {"id": file["id"], "name": file.get("name")},
            "count": len(response.get("permissions", [])),
            "permissions": response.get("permissions", []),
            "next_page_token": response.get("nextPageToken"),
            "has_more": bool(response.get("nextPageToken")),
        }

    def get_permission(self, file_id: str, permission_id: str) -> dict[str, Any]:
        return (
            self.drive.permissions()
            .get(
                fileId=file_id,
                permissionId=permission_id,
                fields="id,type,role,emailAddress,domain,displayName,expirationTime,deleted,allowFileDiscovery",
                supportsAllDrives=True,
            )
            .execute()
        )

    def create_permission(
        self,
        file_id: str,
        *,
        permission_type: str,
        role: str,
        email_address: str | None = None,
        domain: str | None = None,
        allow_file_discovery: bool | None = None,
        send_notification_email: bool = True,
    ) -> dict[str, Any]:
        body: dict[str, Any] = {"type": permission_type, "role": role}
        if email_address:
            body["emailAddress"] = email_address
        if domain:
            body["domain"] = domain
        if allow_file_discovery is not None:
            body["allowFileDiscovery"] = allow_file_discovery
        return (
            self.drive.permissions()
            .create(
                fileId=file_id,
                body=body,
                sendNotificationEmail=send_notification_email
                if permission_type in {"user", "group"}
                else False,
                fields="id,type,role,emailAddress,domain,displayName,expirationTime,allowFileDiscovery",
                supportsAllDrives=True,
            )
            .execute()
        )

    def update_permission(
        self, file_id: str, permission_id: str, role: str
    ) -> dict[str, Any]:
        return (
            self.drive.permissions()
            .update(
                fileId=file_id,
                permissionId=permission_id,
                body={"role": role},
                fields="id,type,role,emailAddress,domain,displayName,expirationTime,allowFileDiscovery",
                supportsAllDrives=True,
            )
            .execute()
        )

    def delete_permission(self, file_id: str, permission_id: str) -> dict[str, Any]:
        self.drive.permissions().delete(
            fileId=file_id, permissionId=permission_id, supportsAllDrives=True
        ).execute()
        return {"file_id": file_id, "permission_id": permission_id, "deleted": True}

    def list_shared_drives(
        self, *, page_size: int = 100, page_token: str | None = None
    ) -> dict[str, Any]:
        params: dict[str, Any] = {
            "pageSize": max(1, min(int(page_size), 100)),
            "fields": "nextPageToken,drives(id,name,createdTime,hidden,restrictions,capabilities)",
        }
        if page_token:
            params["pageToken"] = page_token
        response = self.drive.drives().list(**params).execute()
        return {
            "count": len(response.get("drives", [])),
            "drives": response.get("drives", []),
            "next_page_token": response.get("nextPageToken"),
            "has_more": bool(response.get("nextPageToken")),
        }

    def get_start_page_token(self, *, drive_id: str | None = None) -> dict[str, Any]:
        params: dict[str, Any] = {"supportsAllDrives": True}
        if drive_id:
            params["driveId"] = drive_id
        response = self.drive.changes().getStartPageToken(**params).execute()
        return {"drive_id": drive_id, "start_page_token": response["startPageToken"]}

    def list_changes(
        self,
        page_token: str,
        *,
        drive_id: str | None = None,
        page_size: int = 100,
        include_removed: bool = True,
    ) -> dict[str, Any]:
        params: dict[str, Any] = {
            "pageToken": page_token,
            "pageSize": max(1, min(int(page_size), 1000)),
            "includeRemoved": include_removed,
            "supportsAllDrives": True,
            "fields": "nextPageToken,newStartPageToken,changes(fileId,removed,time,changeType,driveId,file(id,name,mimeType,parents,trashed,modifiedTime,webViewLink))",
        }
        if drive_id:
            params.update({"driveId": drive_id, "includeItemsFromAllDrives": True})
        response = self.drive.changes().list(**params).execute()
        return {
            "drive_id": drive_id,
            "count": len(response.get("changes", [])),
            "changes": response.get("changes", []),
            "next_page_token": response.get("nextPageToken"),
            "new_start_page_token": response.get("newStartPageToken"),
            "has_more": bool(response.get("nextPageToken")),
        }

    def find_child_folder(self, parent_id: str, name: str) -> dict[str, Any] | None:
        safe = self._escape(name)
        files, _ = self._execute_files_page(
            max_items=1,
            q=f"'{parent_id}' in parents and name = '{safe}' and mimeType = '{GOOGLE_FOLDER_MIME}' and trashed = false",
            pageSize=1,
            fields="files(id,name,mimeType,parents)",
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            corpora="allDrives",
        )
        return files[0] if files else None
