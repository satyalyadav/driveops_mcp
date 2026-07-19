"""Content download and local extraction helpers."""

from __future__ import annotations

import io
import mimetypes
import shutil
import stat
import zipfile
from pathlib import Path
from typing import Any

from googleapiclient.http import MediaIoBaseDownload

from .errors import DriveOpsError
from .schemas import GOOGLE_DOC_MIME, GOOGLE_SHEET_MIME, GOOGLE_SLIDE_MIME


def extract_content(
    meta: dict[str, Any],
    data: bytes,
    *,
    exported: bool,
    output_dir: str | None,
    max_text_chars: int,
    overwrite: bool,
) -> dict[str, Any]:
    """Extract supported content from bytes already fetched from Drive."""

    mime = meta.get("mimeType", "")
    name = (meta.get("name") or "").lower()
    result: dict[str, Any] = {
        "id": meta["id"],
        "name": meta.get("name"),
        "mimeType": mime,
        "source_size": len(data),
        "webViewLink": meta.get("webViewLink"),
    }
    if (
        mime.startswith("text/")
        or mime in {"application/json", "application/xml"}
        or exported
    ):
        text = data.decode("utf-8", errors="replace")
    elif mime == "application/pdf" or name.endswith(".pdf"):
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
        result["pages"] = len(reader.pages)
    elif (
        mime
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or name.endswith(".docx")
    ):
        from docx import Document

        document = Document(io.BytesIO(data))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    elif (
        mime
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or name.endswith(".pptx")
    ):
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        slide_text = []
        for number, slide in enumerate(presentation.slides, start=1):
            parts = [
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            ]
            slide_text.append(f"Slide {number}\n" + "\n".join(parts))
        text = "\n\n".join(slide_text)
        result["slides"] = len(presentation.slides)
    elif (
        mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or name.endswith(".xlsx")
    ):
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets = []
        for sheet in workbook.worksheets:
            rows = [
                "\t".join("" if value is None else str(value) for value in row)
                for row in sheet.iter_rows(values_only=True)
            ]
            sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        text = "\n\n".join(sheets)
        result["sheets"] = workbook.sheetnames
    elif mime in {
        "application/zip",
        "application/x-zip-compressed",
    } or name.endswith(".zip"):
        return extract_zip(meta, data, output_dir=output_dir, overwrite=overwrite)
    else:
        raise DriveOpsError(
            "Unsupported extraction format. Supported: Google Docs/Sheets/Slides, text, PDF, DOCX, PPTX, XLSX, and ZIP."
        )

    max_text_chars = max(1, min(int(max_text_chars), 2_000_000))
    result.update(
        {
            "contentType": "text",
            "text": text[:max_text_chars],
            "truncated": len(text) > max_text_chars,
            "characters": len(text),
        }
    )
    return result


def extract_zip(
    meta: dict[str, Any], data: bytes, *, output_dir: str | None, overwrite: bool
) -> dict[str, Any]:
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        entries = archive.infolist()
        if len(entries) > 1000:
            raise DriveOpsError("ZIP archive contains more than 1000 entries.")
        expanded_size = sum(entry.file_size for entry in entries)
        if expanded_size > 100_000_000:
            raise DriveOpsError("ZIP archive expands beyond the 100 MB safety limit.")
        result: dict[str, Any] = {
            "id": meta["id"],
            "name": meta.get("name"),
            "mimeType": meta.get("mimeType"),
            "contentType": "archive",
            "entry_count": len(entries),
            "expanded_size": expanded_size,
            "entries": [
                {
                    "name": entry.filename,
                    "size": entry.file_size,
                    "is_dir": entry.is_dir(),
                }
                for entry in entries
            ],
        }
        if not output_dir:
            result["message"] = (
                "Archive inspected but not unpacked; provide output_dir to extract it."
            )
            return result

        root = Path(output_dir).expanduser().resolve()
        root.mkdir(parents=True, exist_ok=True)
        destinations: list[tuple[zipfile.ZipInfo, Path]] = []
        for entry in entries:
            mode = entry.external_attr >> 16
            if stat.S_ISLNK(mode):
                raise DriveOpsError(
                    f"Refusing to extract ZIP symlink: {entry.filename}"
                )
            destination = (root / entry.filename).resolve()
            if not destination.is_relative_to(root):
                raise DriveOpsError(f"Refusing unsafe ZIP path: {entry.filename}")
            if not entry.is_dir() and destination.exists() and not overwrite:
                raise DriveOpsError(
                    f"Extraction destination already exists: {destination}"
                )
            destinations.append((entry, destination))
        extracted = []
        for entry, destination in destinations:
            if entry.is_dir():
                destination.mkdir(parents=True, exist_ok=True)
                continue
            destination.parent.mkdir(parents=True, exist_ok=True)
            with archive.open(entry) as source, destination.open("wb") as target:
                shutil.copyfileobj(source, target)
            extracted.append(str(destination))
        result.update({"output_dir": str(root), "extracted_paths": extracted})
        return result


def default_export_format(mime_type: str) -> str:
    formats = {
        GOOGLE_DOC_MIME: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        GOOGLE_SHEET_MIME: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        GOOGLE_SLIDE_MIME: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.google-apps.drawing": "application/pdf",
    }
    if mime_type not in formats:
        raise DriveOpsError(
            f"No default export format is known for {mime_type}; provide export_format."
        )
    return formats[mime_type]


def download_name(name: str, content_type: str) -> str:
    suffix = mimetypes.guess_extension(content_type) or ""
    if suffix == ".jpe":
        suffix = ".jpg"
    if suffix and not name.lower().endswith(suffix.lower()):
        return f"{name}{suffix}"
    return name


def download_request(request: Any, max_bytes: int) -> bytes:
    output = io.BytesIO()
    max_bytes = max(1, int(max_bytes))
    downloader = MediaIoBaseDownload(
        output,
        request,
        chunksize=min(max_bytes, 1024 * 1024),
    )
    done = False
    while not done and output.tell() < max_bytes:
        _, done = downloader.next_chunk()
    return output.getvalue()


def content_response(
    meta: dict[str, Any], data: bytes, content_type: str, max_text_bytes: int
) -> dict[str, Any]:
    truncated = len(data) > max_text_bytes
    raw = data[:max_text_bytes]
    return {
        "id": meta.get("id"),
        "name": meta.get("name"),
        "mimeType": meta.get("mimeType"),
        "contentType": content_type,
        "text": raw.decode("utf-8", errors="replace"),
        "truncated": truncated,
        "webViewLink": meta.get("webViewLink"),
        "webContentLink": meta.get("webContentLink"),
    }
